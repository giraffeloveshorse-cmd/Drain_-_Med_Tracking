#!/usr/bin/env python3
"""Generate a PowerPoint presentation for Tracie's Recovery Tracker."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Colors (derived from the app's CSS variables)
# ---------------------------------------------------------------------------
ACCENT_BLUE = RGBColor(0x1D, 0x9B, 0xF0)
DARK_TEXT = RGBColor(0x0F, 0x14, 0x19)
BODY_TEXT = RGBColor(0x1A, 0x1F, 0x26)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
DIM_GRAY = RGBColor(0x71, 0x76, 0x7B)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FLUID_COLORS = {
    "Dark Red": RGBColor(0x8B, 0x1A, 0x1A),
    "Bright Red": RGBColor(0xDC, 0x26, 0x26),
    "Pink": RGBColor(0xF4, 0x72, 0xB6),
    "Amber": RGBColor(0xD9, 0x77, 0x06),
    "Yellow": RGBColor(0xEA, 0xB3, 0x08),
    "Clear": RGBColor(0xA5, 0xF3, 0xFC),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Calibri"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _add_accent_bar(slide):
    """Add a thin blue accent bar across the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def _add_footer(slide, slide_num):
    """Add footer text and slide number."""
    # App name – bottom left
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Tracie's Recovery Tracker"
    p.font.size = Pt(10)
    p.font.color.rgb = DIM_GRAY
    p.font.name = FONT_NAME

    # Slide number – bottom right
    txBox2 = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(0.8), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(slide_num)
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(10)
    p2.font.color.rgb = DIM_GRAY
    p2.font.name = FONT_NAME


def _set_font(run, size=16, bold=False, color=None, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    if color:
        run.font.color.rgb = color


def _add_title(slide, text, left=0.8, top=0.4, width=11.7, size=36):
    """Add a styled title textbox and return the textframe."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=True, color=DARK_TEXT)
    return tf


def _add_subtitle(slide, text, left=0.8, top=1.1, width=11.7, size=18):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, color=DIM_GRAY)
    return tf


def _add_bullet_list(slide, items, left=0.8, top=1.8, width=11.0, font_size=16):
    """Add a bullet-point list. Each item is a string."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"\u2022  {item}"
        _set_font(run, size=font_size, color=BODY_TEXT)
    return tf


def _add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                       font_size=14, font_color=DARK_TEXT, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _set_font(run, size=font_size, bold=bold, color=font_color)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def _add_arrow(slide, start_left, start_top, end_left, end_top):
    """Add a simple connector arrow shape."""
    # Use a line shape
    connector = slide.shapes.add_connector(
        1,  # straight connector
        Inches(start_left), Inches(start_top),
        Inches(end_left), Inches(end_top),
    )
    connector.line.color.rgb = DIM_GRAY
    connector.line.width = Pt(1.5)
    # End arrow
    connector.end_style = "arrow"
    return connector


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_accent_bar(slide)

    # App icon
    icon_path = os.path.join(os.path.dirname(__file__), "icon-512.png")
    if os.path.exists(icon_path):
        slide.shapes.add_picture(icon_path, Inches(5.9), Inches(1.0), Inches(1.5), Inches(1.5))

    # Title
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.8), Inches(9.3), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Tracie's Recovery Tracker"
    _set_font(run, size=44, bold=True, color=DARK_TEXT)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(3.8), Inches(9.3), Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "A Privacy-First PWA for Post-Surgical Recovery Monitoring"
    _set_font(run2, size=20, color=DIM_GRAY)

    # Institution
    txBox3 = slide.shapes.add_textbox(Inches(2.0), Inches(4.6), Inches(9.3), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Lankenau Medical Center"
    _set_font(run3, size=16, color=ACCENT_BLUE, bold=True)

    _add_footer(slide, 1)


def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)
    _add_title(slide, "The Challenge")

    bullets = [
        "Post-surgical patients must track drain output volumes, fluid color changes, "
        "and multiple medications with different dosing intervals",
        "Paper logs are easy to lose, hard to share with providers, "
        "and lack timing reminders",
        "A dedicated native app is overkill for a temporary recovery period",
    ]
    _add_bullet_list(slide, bullets, top=1.6, font_size=18)

    # Goal callout box
    _add_rounded_rect(
        slide, 1.5, 4.8, 10.3, 1.2, LIGHT_BG,
        text='Goal: A zero-install, offline-capable tool that works '
             'immediately on any phone',
        font_size=18, font_color=ACCENT_BLUE, bold=True,
    )

    _add_footer(slide, 2)


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)
    _add_title(slide, "Single-File PWA Architecture")

    # Left column – bullet list
    bullets = [
        "Single index.html file — no build step",
        "React 18 loaded via CDN",
        "Service Worker for offline caching",
        "localStorage for all data (privacy-first)",
        "Installable on home screen (iOS / Android)",
        "Free hosting on GitHub Pages",
    ]
    _add_bullet_list(slide, bullets, left=0.8, top=1.6, width=5.5, font_size=16)

    # Right column – architecture diagram (boxes + arrows)
    cx = 9.5  # center x of diagram column
    # index.html box
    _add_rounded_rect(slide, 7.8, 1.6, 3.4, 1.0, ACCENT_BLUE,
                       text="index.html\nReact 18 + CSS + JS",
                       font_size=14, font_color=WHITE, bold=True)

    # Three boxes below
    _add_rounded_rect(slide, 7.0, 3.5, 2.2, 0.8, LIGHT_BG,
                       text="localStorage\n(Patient Data)", font_size=12, bold=True)
    _add_rounded_rect(slide, 9.5, 3.5, 2.2, 0.8, LIGHT_BG,
                       text="sw.js\n(Offline Cache)", font_size=12, bold=True)
    _add_rounded_rect(slide, 8.25, 5.0, 2.5, 0.8, LIGHT_BG,
                       text="manifest.json\n(PWA Install)", font_size=12, bold=True)

    # Arrows from index.html down
    _add_arrow(slide, 8.8, 2.6, 8.1, 3.5)
    _add_arrow(slide, 10.2, 2.6, 10.6, 3.5)
    _add_arrow(slide, 9.5, 2.6, 9.5, 5.0)

    _add_footer(slide, 3)


def add_drain_tracking_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)
    _add_title(slide, "Drain Tracking")
    _add_subtitle(slide, "3 JP drains  •  Volume in ml  •  Fluid color  •  Timestamped")

    bullets = [
        "Select drain (#1, #2, or #3) with large tap targets",
        "Quick-fill volume buttons: +5, +10, +15, +20, +25, +30 ml",
        "6-color fluid palette maps appearance to healing stage",
        "Optional notes field for each entry",
    ]
    _add_bullet_list(slide, bullets, top=1.8, width=6.0, font_size=16)

    # Fluid color circles
    start_x = 7.5
    y = 2.2
    circle_size = 0.55
    gap = 0.15
    for i, (name, color) in enumerate(FLUID_COLORS.items()):
        x = start_x + i * (circle_size + gap)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(circle_size), Inches(circle_size),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Label below
        lbl = slide.shapes.add_textbox(
            Inches(x - 0.15), Inches(y + circle_size + 0.05),
            Inches(circle_size + 0.3), Inches(0.4),
        )
        p = lbl.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        _set_font(run, size=9, color=DIM_GRAY)

    # Color guide note
    _add_rounded_rect(
        slide, 7.3, 3.5, 5.2, 2.8, LIGHT_BG,
        text=(
            "Color Guide\n\n"
            "Dark Red → Fresh bleeding (normal early)\n"
            "Bright Red → Active bleeding\n"
            "Pink → Improving drainage\n"
            "Amber → Healing in progress\n"
            "Yellow → Near resolution\n"
            "Clear → Fully healed drainage"
        ),
        font_size=11, font_color=BODY_TEXT, bold=False,
    )

    _add_footer(slide, 4)


def add_medication_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)
    _add_title(slide, "Medication Management")
    _add_subtitle(slide, "Interval timers with safety lockouts")

    meds = [
        ("Oxycodone 5mg", "PRN (as needed)", "4-hour interval", "Pain management", ORANGE),
        ("Cefadroxil 500mg", "Scheduled", "12-hour interval", "Antibiotic — 14 days", ACCENT_BLUE),
        ("Senokot", "PRN (as needed)", "12-hour interval", "Stool softener", GREEN),
    ]

    card_width = 3.6
    card_height = 2.8
    start_x = 0.9
    gap = 0.4
    y = 1.8

    for i, (name, med_type, interval, desc, accent) in enumerate(meds):
        x = start_x + i * (card_width + gap)

        # Card background
        _add_rounded_rect(slide, x, y, card_width, card_height, LIGHT_BG)

        # Accent strip at top of card
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(card_width), Inches(0.08),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent
        strip.line.fill.background()

        # Card text
        txBox = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(y + 0.3),
            Inches(card_width - 0.5), Inches(card_height - 0.4),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Med name
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name
        _set_font(run, size=18, bold=True, color=DARK_TEXT)
        p.space_after = Pt(6)

        # Type
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = med_type
        _set_font(run2, size=13, color=accent, bold=True)
        p2.space_after = Pt(4)

        # Interval
        p3 = tf.add_paragraph()
        run3 = p3.add_run()
        run3.text = interval
        _set_font(run3, size=13, color=DIM_GRAY)
        p3.space_after = Pt(4)

        # Description
        p4 = tf.add_paragraph()
        run4 = p4.add_run()
        run4.text = desc
        _set_font(run4, size=13, color=BODY_TEXT)

    # Key features
    features = [
        "Countdown timer shows time until next safe dose",
        "Progress bar fills as interval elapses",
        '"Too Early" lockout prevents accidental double-dosing',
        '"Log Dose Now" activates only when interval has passed',
    ]
    _add_bullet_list(slide, features, left=0.9, top=5.0, width=11.5, font_size=14)

    _add_footer(slide, 5)


def add_history_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)
    _add_title(slide, "History & Data Export")

    # Left column header
    _add_subtitle(slide, "History Tab", left=0.8, top=1.3, size=20)
    left_bullets = [
        "Daily totals broken down by drain",
        "Full chronological log of all entries",
        "Medication dose history",
        "Delete individual entries",
    ]
    _add_bullet_list(slide, left_bullets, left=0.8, top=1.9, width=5.5, font_size=15)

    # Right column header
    _add_subtitle(slide, "Summary & Export", left=7.0, top=1.3, size=20)
    right_bullets = [
        "Today's drain totals with per-drain breakdown",
        "CSV export (copy to clipboard)",
        "JSON backup (copy to clipboard)",
        "Print-formatted medical report for providers",
    ]
    _add_bullet_list(slide, right_bullets, left=7.0, top=1.9, width=5.5, font_size=15)

    # Divider line
    line = slide.shapes.add_connector(
        1, Inches(6.5), Inches(1.5), Inches(6.5), Inches(5.5),
    )
    line.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    line.line.width = Pt(1)

    # Export callout
    _add_rounded_rect(
        slide, 2.0, 5.2, 9.3, 1.0, LIGHT_BG,
        text="All exports designed for easy sharing with healthcare providers",
        font_size=16, font_color=ACCENT_BLUE, bold=True,
    )

    _add_footer(slide, 6)


def add_technical_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)
    _add_title(slide, "Technical Highlights")

    highlights = [
        ("Offline-First", "Service Worker caches all assets; works without internet after first load",
         ACCENT_BLUE),
        ("Installable", "Add to Home Screen on iOS/Android; launches like a native app",
         GREEN),
        ("Privacy", "Zero data leaves the device — no accounts, no analytics, no server",
         ORANGE),
        ("Mobile-Optimized", "Safe-area insets, touch-friendly targets, dark theme for comfort",
         ACCENT_BLUE),
        ("Zero Dependencies", "No npm, no webpack, no build process — single HTML file",
         GREEN),
        ("Free Hosting", "GitHub Pages with instant deployment",
         ORANGE),
    ]

    cols = 2
    card_w = 5.5
    card_h = 1.3
    start_x = 0.9
    start_y = 1.5
    gap_x = 0.6
    gap_y = 0.3

    for i, (title, desc, accent) in enumerate(highlights):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Card
        _add_rounded_rect(slide, x, y, card_w, card_h, LIGHT_BG)

        # Accent dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.35),
            Inches(0.2), Inches(0.2),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(
            Inches(x + 0.6), Inches(y + 0.15),
            Inches(card_w - 0.8), Inches(0.4),
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        _set_font(run, size=16, bold=True, color=DARK_TEXT)

        # Description
        txBox2 = slide.shapes.add_textbox(
            Inches(x + 0.6), Inches(y + 0.6),
            Inches(card_w - 0.8), Inches(0.6),
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = desc
        _set_font(run2, size=13, color=DIM_GRAY)

    _add_footer(slide, 7)


def add_thankyou_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)

    # Icon
    icon_path = os.path.join(os.path.dirname(__file__), "icon-512.png")
    if os.path.exists(icon_path):
        slide.shapes.add_picture(icon_path, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))

    # Thank you
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(3.2), Inches(9.3), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    _set_font(run, size=44, bold=True, color=DARK_TEXT)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(2.0), Inches(4.2), Inches(9.3), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Built with care for Tracie's recovery"
    _set_font(run2, size=20, color=DIM_GRAY)

    # URL placeholder
    txBox3 = slide.shapes.add_textbox(Inches(2.0), Inches(5.2), Inches(9.3), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "github.io/tracies-recovery-tracker"
    _set_font(run3, size=16, color=ACCENT_BLUE)

    _add_footer(slide, 8)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)
    add_problem_slide(prs)
    add_architecture_slide(prs)
    add_drain_tracking_slide(prs)
    add_medication_slide(prs)
    add_history_slide(prs)
    add_technical_slide(prs)
    add_thankyou_slide(prs)

    output_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "Tracies_Recovery_Tracker_Overview.pptx",
    )
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
